"""Generate Swapp architecture PowerPoint for project manager."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT = r"c:\projectflutter\cap-mobile-main\docs\Architecture_Swapp_CapMobile.pptx"

INDIGO = RGBColor(0x46, 0x40, 0xD6)
DARK = RGBColor(0x1B, 0x1D, 0x2B)
GRAY = RGBColor(0x8C, 0x8F, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(slide, text, subtitle=None):
    slide.shapes.title.text = text
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            for p in shape.text_frame.paragraphs:
                p.font.color.rgb = INDIGO
                p.font.bold = True


def add_bullets(slide, items, level0_size=18, level1_size=16):
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sub = item
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = text
            p.level = 0
            p.font.size = Pt(level0_size)
            p.font.color.rgb = DARK
            for s in sub:
                sp = body.add_paragraph()
                sp.text = s
                sp.level = 1
                sp.font.size = Pt(level1_size)
                sp.font.color.rgb = GRAY
        else:
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(level0_size)
            p.font.color.rgb = DARK


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    add_bullets(slide, bullets)
    return slide


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish
    if slide.shapes.title:
        slide.shapes.title.text = title
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.5), Inches(1))
    tf = tx.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = GRAY
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(
        s,
        "Architecture du module Swapp",
        "Cap Mobile — Documentation technique pour chef de projet\n"
        "Flutter · Riverpod · store-api · Août 2026",
    )

    # Slide 2 — Agenda
    add_content_slide(prs, "Agenda", [
        "1. Contexte et objectifs du module Swapp",
        "2. Vue d'ensemble et principes d'architecture",
        "3. Structure des dossiers (lib/swapp + lib/core/apiswap)",
        "4. Glossaire des notions techniques utilisées",
        "5. Couche API : service, modèles JSON, mappers",
        "6. Couche état : Riverpod et providers",
        "7. Couche UI : pages, widgets, modèles vue",
        "8. Flux métier : navigation, scan, fiche produit",
        "9. Fonctionnalité exemple : Colis dépôt",
        "10. Intégrations, avancement et points d'attention",
    ])

    add_section_slide(prs, "Partie 1", "Contexte & vue d'ensemble")

    add_content_slide(prs, "Contexte — Qu'est-ce que Swapp ?", [
        "Swapp = module métier magasin intégré dans l'application Cap Mobile.",
        "Public cible : vendeurs et équipes magasin (terminaux Zebra TC53E).",
        "Objectifs fonctionnels :",
        ("Consultation stock", [
            "Stock magasin par taille (Dispo, Transit, Picking, Dépôt…)",
            "Stock web SFS, stock alentours, avis produit",
        ]),
        ("Ventes & clients", [
            "Panier, clients, catalogue, commandes (en cours de branchement)",
        ]),
        ("Outils terrain", [
            "Scan code-barres / QR, sélection magasin, colis dépôt",
        ]),
        "Backend : API REST store-api (Chaussea) + API Article Cap Mobile existante.",
    ])

    add_content_slide(prs, "Point d'entrée dans Cap Mobile", [
        "Écran d'accueil Cap Mobile → tuile « Swapp »",
        "Fichier : lib/features/home/pages/home_page.dart",
        "Action : Navigator.push(context, SwappMenuPage.fadeRoute())",
        "SwappMenuPage = hub menu (Stock, Ventes, Équipe)",
        "Tuile « Infos produit » → DetailProduitPage2 (fiche produit v2)",
        "Les autres tuiles affichent « bientôt disponible » (SnackBar)",
    ])

    add_content_slide(prs, "Principe architectural — 3 couches", [
        "Séparation stricte des responsabilités (maintenabilité + testabilité) :",
        ("Couche UI — lib/swapp/", [
            "Pages (écrans), Widgets (composants), Models vue (ProductStockView)",
            "Affiche les données, gère interactions utilisateur (tap, scan)",
        ]),
        ("Couche État — Riverpod (providers)", [
            "lib/core/apiswap/*_provider.dart",
            "Centralise loading, erreurs, produit courant, stock web, alentours",
        ]),
        ("Couche API — lib/core/apiswap/", [
            "SwappApiService (HTTP), modèles JSON, mappers API→UI",
            "Seul point de contact avec store-api",
        ]),
        "Règle d'or : l'UI ne parse jamais le JSON brut de l'API.",
    ])

    add_section_slide(prs, "Partie 2", "Structure du code")

    add_content_slide(prs, "Arborescence lib/swapp/", [
        ("pages/", [
            "swapp_menu_page.dart — Menu principal SWApp",
            "detail_produit_page2.dart — Fiche produit v2 (référence actuelle)",
            "detail_produit_page.dart — Fiche produit v1 (legacy)",
            "client_search_page.dart, client_add_page.dart, reserve_produit_page.dart",
        ]),
        ("widgets/", [
            "Composants réutilisables : ReassortChip, ProductPhotoCircle, ColisDepotPackButton",
            "StoreSelectButton, ArticleCodeDialog, QrCameraScannerPage…",
        ]),
        ("models/", [
            "product_stock_view.dart — modèle UI central",
            "stock_column.dart — définition colonnes tableau stock",
        ]),
        ("utils/", [
            "swapp_scan_flow.dart — DataWedge Zebra + UI chargement",
            "swapp_scan_handler.dart — résolution scan → API",
        ]),
        ("data/", ["test_magasins.dart — données magasins de démonstration"]),
    ])

    add_content_slide(prs, "Arborescence lib/core/apiswap/", [
        ("swapp_api_service.dart", [
            "Client HTTP : fetchModeleGlobal, fetchStockWeb, fetchNearbyStock, fetchProductReviews",
        ]),
        ("swapp_api_constants.dart", [
            "URLs endpoints, token Bearer, codeMag défaut, URL photos produit",
        ]),
        ("models/", [
            "modele_global_model.dart, stock_web_item.dart, nearby_stock_item.dart…",
            "Objets miroir de la réponse JSON API",
        ]),
        ("*_mapper.dart", [
            "ProductStockMapper, StockWebMapper, NearbyStockMapper",
            "Transformation JSON/API → objets UI",
        ]),
        ("*_provider.dart", [
            "swappProductProvider, stockWebProvider, nearbyStockProvider, productReviewProvider",
        ]),
    ])

    add_section_slide(prs, "Partie 3", "Glossaire — Notions techniques")

    add_content_slide(prs, "Flutter — Notions UI", [
        ("Widget", [
            "Brique de base de l'interface Flutter (texte, bouton, liste…)",
            "Tout ce que l'utilisateur voit est un arbre de Widgets",
        ]),
        ("StatelessWidget / StatefulWidget", [
            "Stateless = affichage sans état interne (ex. tuile menu)",
            "Stateful = état local (ex. animation, onglet sélectionné)",
        ]),
        ("ConsumerWidget / ConsumerStatefulWidget", [
            "Widgets Riverpod : peuvent lire/écouter les providers",
            "Utilisés sur DetailProduitPage2, SwappMenuPage",
        ]),
        ("Navigator", [
            "Gère la pile de navigation entre écrans (push/pop)",
            "PageRouteBuilder + FadeTransition pour transitions fluides",
        ]),
        ("LayoutBuilder / Wrap / ListView", [
            "LayoutBuilder : calcule largeur disponible (grille 3 colonnes menu)",
            "Wrap : tuiles sans hauteur forcée (évite espaces vides)",
        ]),
    ])

    add_content_slide(prs, "Riverpod — Gestion d'état", [
        ("Provider", [
            "Conteneur de dépendance : expose un service (ex. SwappApiService)",
        ]),
        ("StateNotifier + StateNotifierProvider", [
            "Gère un état mutable de façon contrôlée (ex. SwappProductState)",
            "Méthodes : fetchModele(), copyWith() pour mises à jour immuables",
        ]),
        ("ref.watch(provider)", [
            "Écoute un provider : rebuild automatique de l'UI si l'état change",
        ]),
        ("ref.read(provider.notifier)", [
            "Appelle une action sans écouter (ex. déclencher un fetch API)",
        ]),
        ("ProviderScope", [
            "Racine Riverpod dans main.dart — rend les providers disponibles à toute l'app",
        ]),
    ])

    add_content_slide(prs, "API REST & JSON", [
        ("REST / HTTP GET", [
            "Communication client-serveur via URLs (store-api Chaussea)",
            "SwappApiService utilise le package http de Dart",
        ]),
        ("JSON", [
            "Format d'échange API → { \"codeModele\": \"...\", \"listePrix\": [...] }",
            "Décodage : jsonDecode() puis fromJson() sur les modèles",
        ]),
        ("Bearer Token", [
            "Header Authorization: Bearer <JWT> pour authentifier les appels API",
            "Actuellement dans SwappApiConstants (à sécuriser en production)",
        ]),
        ("Modèle JSON vs Modèle UI", [
            "JSON model = structure exacte API (ModeleGlobalModel)",
            "UI model = structure optimisée écran (ProductStockView)",
            "Mapper = pont entre les deux",
        ]),
    ])

    add_content_slide(prs, "Patterns & concepts métier", [
        ("Mapper (ProductStockMapper)", [
            "Transforme listePrix API → stockBySize pour le tableau UI",
            "Parse libProduit en référence / colorway / taille",
        ]),
        ("Immutabilité", [
            "ProductStockView est immuable : toute mise à jour crée un nouvel objet",
            "Évite bugs de state partagé entre widgets",
        ]),
        ("DataWedge (Zebra)", [
            "Service scanner hardware sur terminaux magasin",
            "SwappDataWedgeListener écoute les scans sans ouvrir la caméra",
        ]),
        ("i18n (internationalisation)", [
            "AppLocalizations : libellés FR / EN / NL",
            "LanguageMenuButton sur fiche produit",
        ]),
        ("Widget Previews (@Preview)", [
            "Flutter 3.38+ : prévisualisation composants dans l'IDE sans lancer l'app",
            "lib/widget_previews/",
        ]),
    ])

    add_section_slide(prs, "Partie 4", "Couche API en détail")

    add_content_slide(prs, "SwappApiService — Endpoints principaux", [
        ("fetchModeleGlobal(codeModele, codeMag)", [
            "Fiche produit complète + stock magasin par taille (listePrix)",
            "Inclut colisDepot, colisDepotResa, prix, photo, réassort",
        ]),
        ("fetchStockWeb(codeModele)", [
            "Stock entrepôt SFS — onglet « Stock Web »",
        ]),
        ("fetchNearbyStock(gencode, codeMag)", [
            "Disponibilité dans magasins alentours — onglet « Alentours »",
        ]),
        ("fetchProductReviews(...)", [
            "Avis clients produit — onglet « Avis »",
        ]),
        "Gestion erreurs : 401 token, 404 modèle, timeout réseau, SocketException",
    ])

    add_content_slide(prs, "ModeleGlobalModel — Champs clés API", [
        "codeModele, libProduit, forme, prixVente, resteALiver (réassort)",
        "libFamille, libTheme, libRayon, libSaison, libTaille, libPlusProduit",
        ("listePrix[] — stock par taille", [
            "taille, stockMag (dispo), stockTransit, stockPicking",
            "stockDepot, stockResa, stockPreResa, stockVol, stockNonVendable…",
        ]),
        ("colisDepotResa[] / colisDepot[]", [
            "Colis dépôt : déblocage résa ou pari colis",
            "sColisPre, nPCBColis, sDesc → libellé chip UI",
        ]),
    ])

    add_content_slide(prs, "ProductStockMapper → ProductStockView", [
        "Entrée : ModeleGlobalModel (JSON API)",
        "Sortie : ProductStockView (objet UI unique pour toute la fiche produit)",
        ("Transformations principales", [
            "libProduit → reference + colorway + size",
            "listePrix → Map<taille, Map<rubrique, quantité>>",
            "colisDepot* → List<ColisDepotChip>",
            "prixVente → price, resteALiver → reassortOk",
            "URL photo : SwappApiConstants.productPhotoUrl(codeModele)",
        ]),
        "Si API colis vide → chips démo (environnement dev uniquement)",
    ])

    add_section_slide(prs, "Partie 5", "État & Providers")

    add_content_slide(prs, "Providers Riverpod Swapp", [
        ("swappProductProvider", [
            "État : product (ProductStockView?), isLoading, error, lastScannedGencode",
            "Action : fetchModele() après scan ou saisie code",
        ]),
        ("stockWebProvider", [
            "Données stock web SFS pour le produit courant",
        ]),
        ("nearbyStockProvider", [
            "Stock magasins alentours (nécessite gencode article)",
        ]),
        ("productReviewProvider", [
            "Liste avis produit pour l'onglet Avis",
        ]),
        ("swappApiServiceProvider", [
            "Injection du service HTTP (une instance partagée)",
        ]),
        "articleProvider (Cap Mobile existant) : résout gencode → codeMod",
    ])

    add_content_slide(prs, "SwappProductState — Cycle de vie", [
        "1. Utilisateur scanne ou saisit un code",
        "2. isLoading = true, error = null",
        "3. Appel API fetchModeleGlobal",
        "4. Succès → ProductStockMapper → product mis à jour, isLoading = false",
        "5. Erreur → message error affiché (_ErrorBanner + SnackBar)",
        "6. UI (DetailProduitPage2) rebuild via ref.watch(swappProductProvider)",
        "copyWith() permet mises à jour partielles sans muter l'état directement",
    ])

    add_section_slide(prs, "Partie 6", "Couche UI & navigation")

    add_content_slide(prs, "SwappMenuPage — Hub navigation", [
        "ConsumerWidget : lit authProvider pour avatar agent et magasin",
        "En-tête blanc : photo collaborateur + titre SWAPP + nom magasin",
        "3 sections en ListView :",
        ("Stock & inventaire (6 tuiles)", [
            "Infos produit ✓, Mvts Stock, REBs, Inventaire, Réceptions, Transferts",
        ]),
        ("Ventes & clients (6 tuiles)", ["Panier, Mon stock, Clients, Catalogue…"]),
        ("Équipe & outils (3 tuiles)", ["e-Resas, My Team, Calculette"]),
        "Layout : Wrap 3 colonnes + constantes _SwappMenuLayout pour espacements",
    ])

    add_content_slide(prs, "DetailProduitPage2 — Fiche produit v2", [
        "ConsumerStatefulWidget — écran principal métier Swapp",
        ("Composants UI", [
            "_Header2 : retour, langue, clients, panier",
            "_ProductCard2 : hero produit (réassort, photo, pills, colis, magasin)",
            "_OverviewGrid2 : cartes récap Dispo / Transit / Picking / Dépôt",
            "_SizeRow2 : lignes stock colorées par taille",
            "_BottomNav2 : 5 onglets navigation bas",
        ]),
        "Scan : SwappDataWedgeListener + QrCameraScannerPage",
        "Magasin : StorePickerDialog + codeMag collaborateur",
    ])

    add_content_slide(prs, "5 onglets fiche produit", [
        ("0 — Stock Mag", ["swappProductProvider — stock magasin par taille"]),
        ("1 — Stock Web", ["stockWebProvider — disponibilité entrepôt SFS"]),
        ("2 — Alentours", ["nearbyStockProvider — autres magasins"]),
        ("3 — Avis", ["productReviewProvider — notes et commentaires"]),
        ("4 — Réserve", ["reserve_mock_data — partiel / mock"]),
        "StockColumns : colonnes dynamiques (masquées si toutes valeurs = 0)",
        "Toggle Mag/Web change la source et les colonnes visibles",
    ])

    add_section_slide(prs, "Partie 7", "Flux métier — Scan produit")

    add_content_slide(prs, "Chaîne complète après un scan", [
        "1. Scan hardware (DataWedge) OU caméra QR (QrCameraScannerPage)",
        "2. processSwappProductScanUi() — dialog « Chargement… » + haptic",
        "3. handleSwappProductScan() :",
        ("   a. articleProvider.fetchArticle(code)", [
            "Résout gencode → codeMod (API Article Cap Mobile)",
        ]),
        ("   b. swappProductProvider.fetchModele(codeMod, codeMag)", [
            "Charge fiche produit store-api",
        ]),
        ("   c. stockWebProvider + nearbyStockProvider", [
            "Pré-charge données onglets Web et Alentours",
        ]),
        "4. UI rebuild : hero + tableau + overview mis à jour",
        "5. Erreur → SnackBar avec message utilisateur",
    ])

    add_content_slide(prs, "Fichiers scan — Rôles", [
        ("swapp_scan_flow.dart", [
            "SwappDataWedgeListener : init/pause/resume DataWedge",
            "processSwappProductScanUi : orchestration UI (dialog, SnackBar)",
        ]),
        ("swapp_scan_handler.dart", [
            "handleSwappProductScan : logique métier sans widget",
            "swappLooksLikeGencode() : détecte EAN 13+ chiffres",
        ]),
        "Pause DataWedge pendant scan caméra → évite double traitement",
    ])

    add_section_slide(prs, "Partie 8", "Exemple : Colis dépôt")

    add_content_slide(prs, "Fonctionnalité Colis dépôt — Bout en bout", [
        ("API", [
            "ModeleGlobalModel.colisDepotResa + colisDepot",
        ]),
        ("Mapper", [
            "ProductStockMapper._buildColisDepotChips()",
            "Label : sColisPre + « Je débloque/pari… » + nPCBColis + sDesc",
        ]),
        ("UI", [
            "ColisDepotPackButton : bouton animé (onglet Stock Mag uniquement)",
            "showColisDepotPopup : popup Wrap + chips cliquables",
            "Confirmation → succès → chip retiré de la session",
        ]),
        "Fichier : lib/swapp/widgets/colis_depot_chips.dart",
        "État session : _releasedColisIds dans DetailProduitPage2",
    ])

    add_section_slide(prs, "Partie 9", "Intégrations & avancement")

    add_content_slide(prs, "Intégrations transverses Cap Mobile", [
        ("authProvider", [
            "Photo agent, nom, magasin courant (SwappMenuPage, StorePicker)",
        ]),
        ("articleProvider", [
            "Résolution article depuis gencode scanné",
        ]),
        ("dataWedgeServiceProvider", [
            "Scanner Zebra TC53E — service core Cap Mobile",
        ]),
        ("AppLocalizations + AppColors", [
            "Traductions FR/EN/NL, thème visuel cohérent Cap Mobile",
        ]),
        "Swapp réutilise l'infra existante ; apiswap = couche API dédiée Swapp",
    ])

    add_content_slide(prs, "État d'avancement fonctionnel", [
        ("✅ Terminé / actif", [
            "Menu SWApp, fiche produit v2, scan DataWedge + QR",
            "Stock mag, web, alentours, avis (API branchées)",
            "Colis dépôt UI, sélecteur magasin, i18n",
            "Widget Previews debug IDE",
        ]),
        ("🔄 En cours / partiel", [
            "Pages clients (search/add), réservation produit",
            "Tuiles menu non branchées (SnackBar « bientôt »)",
        ]),
        ("⚠️ Points techniques", [
            "Token Bearer hardcodé → migration auth dynamique prod",
            "DetailProduitPage v1 legacy (v2 = référence)",
            "Peu de tests automatisés unitaires Swapp",
        ]),
    ])

    add_content_slide(prs, "Schéma récapitulatif — Flux de données", [
        "store-api (JSON)",
        "    ↓ SwappApiService (HTTP)",
        "    ↓ ModeleGlobalModel.fromJson()",
        "    ↓ ProductStockMapper",
        "    ↓ swappProductProvider (Riverpod)",
        "    ↓ ref.watch()",
        "    ↓ DetailProduitPage2 + Widgets",
        "",
        "Scan : DataWedge → articleProvider → swappProductProvider → UI",
        "",
        "Séparation : modifier l'API n'impacte l'UI que via le mapper.",
    ])

    # Slide finale
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(
        s,
        "Merci — Questions ?",
        "Documentation générée depuis le code Cap Mobile\n"
        "Contact technique : équipe Swapp / H.AMIZIANI\n"
        "Fichiers clés : lib/swapp/ · lib/core/apiswap/",
    )

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
